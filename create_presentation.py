"""
Create Azure Differentiation PowerPoint Presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ========== SLIDE 1: Why Azure is Unique ==========
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "KetzAgenticEcomm: Why Azure?"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 120, 212)  # Azure blue

# Subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "3 Unique Microsoft Capabilities That Change E-Commerce"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(20)
subtitle_para.font.color.rgb = RGBColor(80, 80, 80)

# Content Box
content_top = Inches(1.7)
content_box = slide1.shapes.add_textbox(Inches(0.5), content_top, Inches(9), Inches(5))
text_frame = content_box.text_frame
text_frame.word_wrap = True

# Section 1: GPT-4o Realtime
p1 = text_frame.paragraphs[0]
p1.text = "1. GPT-4o Realtime API"
p1.font.size = Pt(24)
p1.font.bold = True
p1.font.color.rgb = RGBColor(0, 120, 212)
p1.space_after = Pt(6)

p2 = text_frame.add_paragraph()
p2.text = "   • AZURE EXCLUSIVE - Not available on AWS or Google"
p2.font.size = Pt(16)
p2.space_after = Pt(4)

p3 = text_frame.add_paragraph()
p3.text = "   • <300ms voice latency with native barge-in"
p3.font.size = Pt(16)
p3.space_after = Pt(4)

p4 = text_frame.add_paragraph()
p4.text = "   • Built-in function calling from voice commands"
p4.font.size = Pt(16)
p4.space_after = Pt(12)

# Section 2: Unified AI Search
p5 = text_frame.add_paragraph()
p5.text = "2. Unified AI Search Platform"
p5.font.size = Pt(24)
p5.font.bold = True
p5.font.color.rgb = RGBColor(0, 120, 212)
p5.space_after = Pt(6)

p6 = text_frame.add_paragraph()
p6.text = "   • Vector + Semantic + Keyword in ONE index (not on AWS/Google)"
p6.font.size = Pt(16)
p6.space_after = Pt(4)

p7 = text_frame.add_paragraph()
p7.text = "   • Native agentic retrieval - Query decomposition built-in"
p7.font.size = Pt(16)
p7.space_after = Pt(4)

p8 = text_frame.add_paragraph()
p8.text = "   • Live synonym updates without reindexing"
p8.font.size = Pt(16)
p8.space_after = Pt(12)

# Section 3: Enterprise GPT-4o
p9 = text_frame.add_paragraph()
p9.text = "3. Enterprise GPT-4o with Guarantees"
p9.font.size = Pt(24)
p9.font.bold = True
p9.font.color.rgb = RGBColor(0, 120, 212)
p9.space_after = Pt(6)

p10 = text_frame.add_paragraph()
p10.text = "   • Data residency + No training on customer data (contractual)"
p10.font.size = Pt(16)
p10.space_after = Pt(4)

p11 = text_frame.add_paragraph()
p11.text = "   • Latest models: GPT-4o, GPT-4o Vision (AWS/Google have older models)"
p11.font.size = Pt(16)
p11.space_after = Pt(4)

# ========== SLIDE 2: Quick Wins ==========
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame2 = title_box2.text_frame
title_frame2.text = "Azure vs AWS/Google: Quick Comparison"
title_para2 = title_frame2.paragraphs[0]
title_para2.font.size = Pt(40)
title_para2.font.bold = True
title_para2.font.color.rgb = RGBColor(0, 120, 212)

# Left side - Azure
left_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(5.5))
left_frame = left_box.text_frame
left_frame.word_wrap = True

l1 = left_frame.paragraphs[0]
l1.text = "Azure (Unified Platform) ✅"
l1.font.size = Pt(20)
l1.font.bold = True
l1.font.color.rgb = RGBColor(0, 120, 212)
l1.space_after = Pt(10)

l2 = left_frame.add_paragraph()
l2.text = "Voice: GPT-4o Realtime"
l2.font.size = Pt(16)
l2.level = 0
l2.space_after = Pt(6)

l3 = left_frame.add_paragraph()
l3.text = "Native WebRTC, <300ms"
l3.font.size = Pt(14)
l3.level = 1
l3.space_after = Pt(8)

l4 = left_frame.add_paragraph()
l4.text = "Search: AI Search"
l4.font.size = Pt(16)
l4.level = 0
l4.space_after = Pt(6)

l5 = left_frame.add_paragraph()
l5.text = "Vector+Semantic+Keyword"
l5.font.size = Pt(14)
l5.level = 1
l5.space_after = Pt(8)

l6 = left_frame.add_paragraph()
l6.text = "Agentic: Knowledge Bases"
l6.font.size = Pt(16)
l6.level = 0
l6.space_after = Pt(6)

l7 = left_frame.add_paragraph()
l7.text = "Auto query decomposition"
l7.font.size = Pt(14)
l7.level = 1
l7.space_after = Pt(8)

l8 = left_frame.add_paragraph()
l8.text = "Synonyms: Live Updates"
l8.font.size = Pt(16)
l8.level = 0
l8.space_after = Pt(6)

l9 = left_frame.add_paragraph()
l9.text = "No reindexing needed"
l9.font.size = Pt(14)
l9.level = 1

# Right side - AWS/Google
right_box = slide2.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.5))
right_frame = right_box.text_frame
right_frame.word_wrap = True

r1 = right_frame.paragraphs[0]
r1.text = "AWS/Google (Patchwork) ⚠️"
r1.font.size = Pt(20)
r1.font.bold = True
r1.font.color.rgb = RGBColor(200, 80, 0)
r1.space_after = Pt(10)

r2 = right_frame.add_paragraph()
r2.text = "Voice: Lex + Polly + Lambda"
r2.font.size = Pt(16)
r2.level = 0
r2.space_after = Pt(6)

r3 = right_frame.add_paragraph()
r3.text = "Custom orchestration, ~1-2s"
r3.font.size = Pt(14)
r3.level = 1
r3.font.color.rgb = RGBColor(150, 150, 150)
r3.space_after = Pt(8)

r4 = right_frame.add_paragraph()
r4.text = "Search: OpenSearch (separate)"
r4.font.size = Pt(16)
r4.level = 0
r4.space_after = Pt(6)

r5 = right_frame.add_paragraph()
r5.text = "No unified index"
r5.font.size = Pt(14)
r5.level = 1
r5.font.color.rgb = RGBColor(150, 150, 150)
r5.space_after = Pt(8)

r6 = right_frame.add_paragraph()
r6.text = "Agentic: Bedrock Agents"
r6.font.size = Pt(16)
r6.level = 0
r6.space_after = Pt(6)

r7 = right_frame.add_paragraph()
r7.text = "Requires Lambda orchestration"
r7.font.size = Pt(14)
r7.level = 1
r7.font.color.rgb = RGBColor(150, 150, 150)
r7.space_after = Pt(8)

r8 = right_frame.add_paragraph()
r8.text = "Synonyms: Reindex Required"
r8.font.size = Pt(16)
r8.level = 0
r8.space_after = Pt(6)

r9 = right_frame.add_paragraph()
r9.text = "Downtime for updates"
r9.font.size = Pt(14)
r9.level = 1
r9.font.color.rgb = RGBColor(150, 150, 150)

# Bottom message
bottom_box = slide2.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
bottom_frame = bottom_box.text_frame
bottom_para = bottom_frame.paragraphs[0]
bottom_para.text = "Key Takeaway: Azure = One unified platform. AWS/Google = 5-6 services to stitch together."
bottom_para.font.size = Pt(18)
bottom_para.font.bold = True
bottom_para.font.color.rgb = RGBColor(0, 120, 212)
bottom_para.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = "KetzAgenticEcomm_Azure_Differentiation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved: {output_path}")
